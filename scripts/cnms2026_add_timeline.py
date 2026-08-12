"""Add the CNMS 2026 "Timeline" slide layout (+ example overlay slide) to a deck.

Usage:
    python scripts/cnms2026_add_timeline.py [input.pptx] [output.pptx]

With no input, builds a standalone 16:9 deck containing just the Timeline
layout and the example slide — paste the example slide into another deck with
"Keep Source Formatting" and the layout travels with it.

Timeline geometry (16:9, 13.333 x 7.5 in): five equal segments of exactly
slide_width/5 = 2.6667 in; node k (1-based) is centered at x = (k-0.5)*2.6667 in
on the 3 pt spine at y = 6.35 in. Overlay boxes that dim a whole segment span
x = (k-1)*2.6667 in, width 2.6667 in, y = 5.95 in, height 1.55 in.
Colors: theme navy 0E2841 spine, teal 156082 nodes (node 5 orange E97132).
"""
import sys
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT, CONTENT_TYPE as CT
from pptx.opc.packuri import PackURI
from pptx.parts.slide import SlideLayoutPart
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

IN = 914400
NAVY = RGBColor(0x0E, 0x28, 0x41)
TEAL = RGBColor(0x15, 0x60, 0x82)
ORANGE = RGBColor(0xE9, 0x71, 0x32)
GRAY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

LABELS = [
    (["Line completion /", "multi-line edits"], "IntelliCode-style smart complete"),
    (["GitHub Copilot Chat"], None),
    (["GitHub Copilot Cloud Agent"], None),
    (["GitHub Copilot CLI"], None),
    (["Increasingly multi-agent,", "multi-tool workflows"], None),
]


def add_timeline(prs):
    master = prs.slide_masters[0]
    blank = [l for l in master.slide_layouts if l.name == "Blank"][0]

    # clone the Blank layout as a new "Timeline" layout part
    package = prs.part.package
    existing = [str(l.part.partname) for m in prs.slide_masters for l in m.slide_layouts]
    n = 1
    while "/ppt/slideLayouts/slideLayout%d.xml" % n in existing:
        n += 1
    new_part = SlideLayoutPart.load(PackURI("/ppt/slideLayouts/slideLayout%d.xml" % n),
                                    CT.PML_SLIDE_LAYOUT, package, blank.part.blob)
    new_part.relate_to(master.part, RT.SLIDE_MASTER)
    layout_el = new_part._element
    layout_el.attrib.pop("type", None)
    layout_el.set("userDrawn", "1")
    layout_el.find(qn("p:cSld")).set("name", "Timeline")

    rId = master.part.relate_to(new_part, RT.SLIDE_LAYOUT)
    lst = master.element.find(qn("p:sldLayoutIdLst"))
    max_id = max(int(e.get("id")) for m in prs.slide_masters
                 for e in m.element.find(qn("p:sldLayoutIdLst")))
    lid = etree.SubElement(lst, qn("p:sldLayoutId"))
    lid.set("id", str(max_id + 1))
    lid.set(qn("r:id"), rId)
    layout = [l for l in master.slide_layouts if l.name == "Timeline"][0]

    SEG = prs.slide_width // 5
    centers = [k * SEG + SEG // 2 for k in range(5)]
    line_y = int(6.35 * IN)
    diams = [int(d * IN) for d in (0.30, 0.38, 0.46, 0.54, 0.62)]

    # draw on a scratch slide, then move the shapes into the layout's spTree
    scratch = prs.slides.add_slide(layout)
    pre_ids = {sp.shape_id for sp in scratch.shapes}
    shp = scratch.shapes

    line = shp.add_connector(MSO_CONNECTOR.STRAIGHT,
                             Emu(int(0.40 * IN)), Emu(line_y),
                             Emu(int(13.05 * IN)), Emu(line_y))
    line.name = "Timeline Spine"
    line.line.color.rgb = NAVY
    line.line.width = Pt(3)
    ln = line.line._get_or_add_ln()
    ln.set("cap", "rnd")
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle"); tail.set("w", "lg"); tail.set("len", "lg")

    for k, (cx, d, (title_lines, caption)) in enumerate(zip(centers, diams, LABELS), 1):
        node = shp.add_shape(MSO_SHAPE.OVAL, Emu(cx - d // 2), Emu(line_y - d // 2),
                             Emu(d), Emu(d))
        node.name = "Timeline Node %d" % k
        node.fill.solid()
        node.fill.fore_color.rgb = ORANGE if k == 5 else TEAL
        node.line.color.rgb = WHITE
        node.line.width = Pt(1.5)
        node.shadow.inherit = False

        box_w = int(2.55 * IN)
        left = min(max(cx - box_w // 2, 0), prs.slide_width - box_w)
        tb = shp.add_textbox(Emu(left), Emu(int(6.75 * IN)),
                             Emu(box_w), Emu(int(0.68 * IN)))
        tb.name = "Timeline Label %d" % k
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        for i, text in enumerate(title_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = text
            r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = NAVY
        if caption:
            pc = tf.add_paragraph(); pc.alignment = PP_ALIGN.CENTER
            rc = pc.add_run(); rc.text = caption
            rc.font.size = Pt(9); rc.font.italic = True; rc.font.color.rgb = GRAY

    sptree = layout.shapes._spTree
    lay_ids = [int(e.get("id")) for e in sptree.findall(".//" + qn("p:cNvPr"))]
    next_id = max(lay_ids or [1]) + 1
    for sp_el in list(scratch.shapes._spTree):
        if etree.QName(sp_el).localname not in ("sp", "cxnSp", "pic", "graphicFrame", "grpSp"):
            continue
        cnv = sp_el.find(".//" + qn("p:cNvPr"))
        if int(cnv.get("id")) in pre_ids:
            continue
        cnv.set("id", str(next_id)); next_id += 1
        sptree.append(sp_el)

    sldIdLst = prs.slides._sldIdLst
    sldId = list(sldIdLst)[-1]
    prs.part.drop_rel(sldId.get(qn("r:id")))
    sldIdLst.remove(sldId)

    # example slide: white 72%-opacity overlays dim nodes 2-5, node 1 emphasized
    ex = prs.slides.add_slide(layout)
    band_y, band_h = int(5.95 * IN), int(1.55 * IN)
    for k in range(1, 5):
        ov = ex.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(k * SEG), Emu(band_y),
                                 Emu(SEG), Emu(band_h))
        ov.name = "Dim Node %d" % (k + 1)
        ov.fill.solid()
        ov.fill.fore_color.rgb = WHITE
        srgb = ov.fill.fore_color._xFill.find(".//" + qn("a:srgbClr"))
        etree.SubElement(srgb, qn("a:alpha")).set("val", "72000")
        ov.line.fill.background()
        ov.shadow.inherit = False
    note = ex.shapes.add_textbox(Emu(int(0.4 * IN)), Emu(int(0.25 * IN)),
                                 Emu(int(10.5 * IN)), Emu(int(0.4 * IN)))
    note.name = "Example note"
    r = note.text_frame.paragraphs[0].add_run()
    r.text = ("Example: Timeline layout + 72%-white overlays dimming nodes 2-5. "
              "Duplicate, shift overlays, add your GIF. Safe to delete.")
    r.font.size = Pt(10); r.font.italic = True; r.font.color.rgb = GRAY
    return prs


if __name__ == "__main__":
    inp = sys.argv[1] if len(sys.argv) > 1 else None
    out = sys.argv[2] if len(sys.argv) > 2 else "timeline-out.pptx"
    if inp:
        prs = Presentation(inp)
    else:
        prs = Presentation()
        prs.slide_width, prs.slide_height = 12192000, 6858000
    add_timeline(prs)
    prs.save(out)
    print("saved", out)
